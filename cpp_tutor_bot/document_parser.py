"""
文档解析模块 - 支持PDF、PPT、TXT等格式
"""

import os
from typing import List
from abc import ABC, abstractmethod
from langchain_core.documents import Document

from langchain_text_splitters import RecursiveCharacterTextSplitter


class DocumentParser(ABC):
    """文档解析器基类"""
    
    @abstractmethod
    def parse(self, file_path: str) -> str:
        """解析文档并返回文本内容"""
        pass


class TXTParser(DocumentParser):
    """TXT文件解析器"""

    def parse(self, file_path: str) -> str:
        # 尝试多种编码，提高兼容性
        encodings = ['utf-8', 'gbk', 'gb2312', 'gb18030', 'latin-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except (UnicodeDecodeError, UnicodeError):
                continue
        
        # 如果所有编码都失败，使用errors='ignore'
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()


class PDFParser(DocumentParser):
    """PDF文件解析器"""
    
    def parse(self, file_path: str) -> str:
        #from PyPDF2 import PdfReader
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text


class PPTXParser(DocumentParser):
    """PPTX文件解析器"""
    
    def parse(self, file_path: str) -> str:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text += paragraph.text + "\n"
        return text


class DOCXParser(DocumentParser):
    """DOCX文件解析器"""
    
    def parse(self, file_path: str) -> str:
        from docx import Document as DocxDocument
        
        doc = DocxDocument(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text


class DocumentProcessor:
    """文档处理器 - 负责解析和分割文档"""
    
    def __init__(self):
        self.parsers = {
            '.txt': TXTParser(),
            '.pdf': PDFParser(),
            '.pptx': PPTXParser(),
            '.ppt': PPTXParser(),
            '.docx': DOCXParser(),
        }
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            separators=["\n\n", "\n", "。", "！", "？", "；", "，", " ", ""]
        )
    
    def get_supported_extensions(self) -> List[str]:
        """获取支持的文件扩展名"""
        return list(self.parsers.keys())
    
    def parse_file(self, file_path: str) -> str:
        """解析文件并返回文本"""
        ext = os.path.splitext(file_path)[1].lower()
        if ext not in self.parsers:
            raise ValueError(f"不支持的文件格式: {ext}，支持的格式: {self.get_supported_extensions()}")
        
        parser = self.parsers[ext]
        return parser.parse(file_path)
    
    def process_file(self, file_path: str, source_name: str = None) -> List[Document]:
        """处理文件并返回分割后的文档块"""
        text = self.parse_file(file_path)
        source = source_name or os.path.basename(file_path)
        
        # 分割文本
        chunks = self.text_splitter.split_text(text)
        
        # 创建Document对象
        documents = [
            Document(
                page_content=chunk,
                metadata={"source": source, "chunk_index": i}
            )
            for i, chunk in enumerate(chunks)
        ]
        
        return documents
    
    def process_text(self, text: str, source_name: str = "manual_input") -> List[Document]:
        """处理纯文本并返回分割后的文档块"""
        chunks = self.text_splitter.split_text(text)
        
        documents = [
            Document(
                page_content=chunk,
                metadata={"source": source_name, "chunk_index": i}
            )
            for i, chunk in enumerate(chunks)
        ]
        
        return documents
